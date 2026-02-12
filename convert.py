#!/usr/bin/env python3
"""
Markdown to Office Converter
Converts Markdown files to PDF, DOCX, XLSX, and PPTX formats
"""
import argparse
import io
import os
import re
import sys
from typing import List, Optional, Tuple

import markdown2

# PDF support
try:
    from xhtml2pdf import pisa
    from pypdf import PdfReader, PdfWriter
    import pdfplumber
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

# Office format support
try:
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from openpyxl import Workbook, load_workbook
    from openpyxl.styles import Font, PatternFill, Alignment
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.util import Inches as PptxInches, Pt as PptxPt
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


# ============================================================================
# PDF Functions (existing)
# ============================================================================

def get_html_template(content: str) -> str:
    """Generate HTML template with styling for PDF conversion."""
    return f"""
<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <style>
        @page {{
            size: letter landscape;
            margin: 0.5in;
        }}
        body {{
            font-family: Arial, sans-serif;
        }}
        table {{
            border-collapse: collapse;
            width: 100%;
            font-size: 8px;
        }}
        th, td {{
            border: 1px solid #333;
            padding: 6px;
            text-align: left;
        }}
        th {{
            background-color: #4CAF50;
            color: white;
            font-weight: bold;
        }}
        tr:nth-child(even) {{
            background-color: #f2f2f2;
        }}
    </style>
</head>
<body>
    {content}
</body>
</html>
"""


def markdown_to_pdf_bytes(md_content: str) -> bytes:
    """Convert markdown content to PDF bytes."""
    html_content = markdown2.markdown(md_content, extras=['tables'])
    full_html = get_html_template(html_content)
    
    pdf_buffer = io.BytesIO()
    pisa_status = pisa.CreatePDF(full_html, dest=pdf_buffer)
    
    if pisa_status.err:
        raise RuntimeError("Error creating PDF from markdown")
    
    return pdf_buffer.getvalue()


# ============================================================================
# DOCX Functions
# ============================================================================

def markdown_to_docx(md_content: str, output_file: str) -> None:
    """Convert markdown to Word document."""
    doc = Document()
    
    # Parse markdown
    lines = md_content.split('\n')
    i = 0
    
    while i < len(lines):
        line = lines[i].strip()
        
        if not line:
            i += 1
            continue
        
        # Headers
        if line.startswith('# '):
            heading = doc.add_heading(line[2:], level=1)
        elif line.startswith('## '):
            heading = doc.add_heading(line[3:], level=2)
        elif line.startswith('### '):
            heading = doc.add_heading(line[4:], level=3)
        
        # Tables
        elif line.startswith('|'):
            table_lines = [line]
            i += 1
            while i < len(lines) and lines[i].strip().startswith('|'):
                table_lines.append(lines[i].strip())
                i += 1
            i -= 1
            
            # Parse table
            rows = []
            for tline in table_lines:
                if '---' in tline:  # Skip separator line
                    continue
                cells = [cell.strip() for cell in tline.split('|')[1:-1]]
                rows.append(cells)
            
            if rows:
                table = doc.add_table(rows=len(rows), cols=len(rows[0]))
                table.style = 'Light Grid Accent 1'
                
                for row_idx, row_data in enumerate(rows):
                    for col_idx, cell_data in enumerate(row_data):
                        table.rows[row_idx].cells[col_idx].text = cell_data
                        if row_idx == 0:  # Header row
                            cell = table.rows[row_idx].cells[col_idx]
                            cell.paragraphs[0].runs[0].bold = True
        
        # Regular paragraphs
        else:
            # Simple bold/italic support
            para = doc.add_paragraph()
            text = line
            
            # Handle bold **text**
            parts = re.split(r'(\*\*.*?\*\*)', text)
            for part in parts:
                if part.startswith('**') and part.endswith('**'):
                    para.add_run(part[2:-2]).bold = True
                else:
                    para.add_run(part)
        
        i += 1
    
    doc.save(output_file)


# ============================================================================
# XLSX Functions
# ============================================================================

def markdown_to_xlsx(md_content: str, output_file: str) -> None:
    """Convert markdown tables to Excel spreadsheet."""
    wb = Workbook()
    ws = wb.active
    ws.title = "Data"
    
    # Parse markdown to find tables
    lines = md_content.split('\n')
    table_num = 0
    current_row = 1
    
    for i, line in enumerate(lines):
        line = line.strip()
        
        # Add section headers as worksheet sections
        if line.startswith('# ') or line.startswith('## '):
            header_text = line.lstrip('#').strip()
            cell = ws.cell(row=current_row, column=1, value=header_text)
            cell.font = Font(bold=True, size=14)
            current_row += 2
        
        # Parse tables
        elif line.startswith('|'):
            table_lines = [line]
            j = i + 1
            while j < len(lines) and lines[j].strip().startswith('|'):
                table_lines.append(lines[j].strip())
                j += 1
            
            # Parse table
            rows = []
            for tline in table_lines:
                if '---' in tline:  # Skip separator
                    continue
                cells = [cell.strip() for cell in tline.split('|')[1:-1]]
                rows.append(cells)
            
            if rows:
                for row_data in rows:
                    for col_idx, cell_data in enumerate(row_data, start=1):
                        cell = ws.cell(row=current_row, column=col_idx, value=cell_data)
                        
                        # Style header row
                        if current_row == 1 or rows.index(row_data) == 0:
                            cell.font = Font(bold=True, color="FFFFFF")
                            cell.fill = PatternFill(start_color="4CAF50", end_color="4CAF50", fill_type="solid")
                            cell.alignment = Alignment(horizontal="center")
                    
                    current_row += 1
                
                current_row += 2  # Add spacing between tables
    
    # Auto-adjust column widths
    for column in ws.columns:
        max_length = 0
        column_letter = column[0].column_letter
        for cell in column:
            if cell.value:
                max_length = max(max_length, len(str(cell.value)))
        adjusted_width = min(max_length + 2, 50)
        ws.column_dimensions[column_letter].width = adjusted_width
    
    wb.save(output_file)


# ============================================================================
# PPTX Functions
# ============================================================================

def extract_markdown_sections(md_content: str) -> List[Tuple[str, str]]:
    """Extract sections from markdown based on ## headers."""
    sections = []
    lines = md_content.split('\n')
    current_header = None
    current_content = []
    
    for line in lines:
        if line.strip().startswith('## '):
            if current_header is not None:
                sections.append((current_header, '\n'.join(current_content)))
            current_header = line.strip()[3:].strip()
            current_content = []
        else:
            current_content.append(line)
    
    if current_header is not None:
        sections.append((current_header, '\n'.join(current_content)))
    
    return sections


def markdown_to_pptx(md_content: str, output_file: str) -> None:
    """Convert markdown sections to PowerPoint presentation."""
    prs = Presentation()
    prs.slide_width = PptxInches(10)
    prs.slide_height = PptxInches(7.5)
    
    sections = extract_markdown_sections(md_content)
    
    if not sections:
        print("Warning: No ## headers found in markdown. Creating single slide.")
        sections = [("Untitled", md_content)]
    
    for title, content in sections:
        # Add slide with title and content layout
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Set content
        if len(slide.shapes) > 1:
            content_shape = slide.shapes[1]
            if hasattr(content_shape, 'text_frame'):
                tf = content_shape.text_frame
                tf.text = content.strip()
    
    prs.save(output_file)


def append_to_pptx(markdown_file: str, existing_pptx: str, output_file: str) -> None:
    """Append markdown slides to existing presentation."""
    with open(markdown_file, 'r', encoding='utf-8') as f:
        md_content = f.read()
    
    # Load existing presentation
    prs = Presentation(existing_pptx)
    original_count = len(prs.slides)
    
    # Add new slides
    sections = extract_markdown_sections(md_content)
    for title, content in sections:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        
        if len(slide.shapes) > 1 and hasattr(slide.shapes[1], 'text_frame'):
            slide.shapes[1].text_frame.text = content.strip()
    
    prs.save(output_file)
    print(f"Presentation created successfully: {output_file}")
    print(f"  Original slides: {original_count}")
    print(f"  Appended slides: {len(sections)}")
    print(f"  Total slides: {len(prs.slides)}")


def normalize_text(text: str) -> str:
    """Normalize text for comparison."""
    return re.sub(r'\s+', ' ', text.lower().strip())


def replace_pptx_slides(markdown_file: str, existing_pptx: str, output_file: str) -> None:
    """Replace slides in presentation where titles match markdown headers."""
    with open(markdown_file, 'r', encoding='utf-8') as f:
        md_content = f.read()
    
    sections = extract_markdown_sections(md_content)
    if not sections:
        print("Warning: No ## headers found in markdown")
        print("Using append mode instead...")
        append_to_pptx(markdown_file, existing_pptx, output_file)
        return
    
    # Load existing presentation
    prs = Presentation(existing_pptx)
    
    # Build section map
    section_map = {normalize_text(title): (title, content) for title, content in sections}
    
    # Find matches and track replacements
    matches = []
    for idx, slide in enumerate(prs.slides):
        if slide.shapes.title and slide.shapes.title.text:
            slide_title = slide.shapes.title.text
            normalized = normalize_text(slide_title)
            
            if normalized in section_map:
                matches.append((idx, slide_title, section_map[normalized]))
    
    if not matches:
        print("Warning: No matching slide titles found")
        print("Using append mode instead...")
        append_to_pptx(markdown_file, existing_pptx, output_file)
        return
    
    print(f"Found {len(matches)} matching slides to replace:")
    for idx, title, _ in matches:
        print(f"  Slide {idx + 1}: {title}")
    
    # Replace matched slides
    for idx, old_title, (new_title, content) in matches:
        slide = prs.slides[idx]
        
        # Update title
        slide.shapes.title.text = new_title
        
        # Update content
        if len(slide.shapes) > 1 and hasattr(slide.shapes[1], 'text_frame'):
            tf = slide.shapes[1].text_frame
            tf.clear()
            tf.text = content.strip()
    
    prs.save(output_file)
    print(f"\nPresentation created successfully: {output_file}")
    print(f"  Total slides: {len(prs.slides)}")
    print(f"  Slides replaced: {len(matches)}")


# ============================================================================
# PDF append/replace (existing functions)
# ============================================================================

def create_pdf(markdown_file: str, output_file: str) -> None:
    """Create a new PDF from markdown file."""
    with open(markdown_file, 'r', encoding='utf-8') as f:
        md_content = f.read()
    
    pdf_bytes = markdown_to_pdf_bytes(md_content)
    
    with open(output_file, 'wb') as f:
        f.write(pdf_bytes)
    
    print(f"PDF created successfully: {output_file}")


def append_to_pdf(markdown_file: str, existing_pdf: str, output_file: str) -> None:
    """Append markdown content as new pages to an existing PDF."""
    with open(markdown_file, 'r', encoding='utf-8') as f:
        md_content = f.read()
    
    new_pdf_bytes = markdown_to_pdf_bytes(md_content)
    
    reader_existing = PdfReader(existing_pdf)
    reader_new = PdfReader(io.BytesIO(new_pdf_bytes))
    
    writer = PdfWriter()
    
    for page in reader_existing.pages:
        writer.add_page(page)
    
    for page in reader_new.pages:
        writer.add_page(page)
    
    with open(output_file, 'wb') as f:
        writer.write(f)
    
    print(f"PDF created successfully: {output_file}")
    print(f"  Original pages: {len(reader_existing.pages)}")
    print(f"  Appended pages: {len(reader_new.pages)}")
    print(f"  Total pages: {len(reader_existing.pages) + len(reader_new.pages)}")


def extract_page_titles(pdf_path: str) -> List[Tuple[int, str]]:
    """Extract title from each PDF page."""
    if pdfplumber is None:
        return []
    
    page_titles = []
    with pdfplumber.open(pdf_path) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text()
            if text:
                lines = [line.strip() for line in text.split('\n') if line.strip()]
                if lines:
                    page_titles.append((i, lines[0]))
    
    return page_titles


def replace_pdf_pages(markdown_file: str, existing_pdf: str, output_file: str) -> None:
    """Replace PDF pages where markdown headers match page titles."""
    with open(markdown_file, 'r', encoding='utf-8') as f:
        md_content = f.read()
    
    sections = extract_markdown_sections(md_content)
    if not sections:
        print("Warning: No ## headers found in markdown file")
        print("Using append mode instead...")
        append_to_pdf(markdown_file, existing_pdf, output_file)
        return
    
    page_titles = extract_page_titles(existing_pdf)
    section_map = {normalize_text(header): content for header, content in sections}
    
    matches = []
    for page_num, page_title in page_titles:
        normalized = normalize_text(page_title)
        if normalized in section_map:
            matches.append((page_num, page_title, section_map[normalized]))
    
    if not matches:
        print("Warning: No matching headers found between PDF and markdown")
        print("Using append mode instead...")
        append_to_pdf(markdown_file, existing_pdf, output_file)
        return
    
    print(f"Found {len(matches)} matching pages to replace:")
    for page_num, title, _ in matches:
        print(f"  Page {page_num + 1}: {title}")
    
    reader_existing = PdfReader(existing_pdf)
    writer = PdfWriter()
    
    for page_num in range(len(reader_existing.pages)):
        match = next((m for m in matches if m[0] == page_num), None)
        
        if match:
            _, title, section_content = match
            section_pdf_bytes = markdown_to_pdf_bytes(section_content)
            section_reader = PdfReader(io.BytesIO(section_pdf_bytes))
            
            for section_page in section_reader.pages:
                writer.add_page(section_page)
        else:
            writer.add_page(reader_existing.pages[page_num])
    
    with open(output_file, 'wb') as f:
        writer.write(f)
    
    print(f"\nPDF created successfully: {output_file}")
    print(f"  Original pages: {len(reader_existing.pages)}")
    print(f"  Pages replaced: {len(matches)}")


# ============================================================================
# Main CLI
# ============================================================================

def main():
    parser = argparse.ArgumentParser(
        description='Convert Markdown to PDF, DOCX, XLSX, or PPTX formats',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  # Create PDF (default)
  python convert.py input.md
  
  # Create Word document
  python convert.py input.md --format docx -o output.docx
  
  # Create Excel spreadsheet
  python convert.py input.md --format xlsx -o output.xlsx
  
  # Create PowerPoint presentation
  python convert.py input.md --format pptx -o output.pptx
  
  # Append to existing presentation
  python convert.py new_slides.md --format pptx --append existing.pptx -o updated.pptx
  
  # Replace matching slides in presentation
  python convert.py updates.md --format pptx --replace presentation.pptx -o final.pptx
        """
    )
    
    parser.add_argument('markdown_file', help='Input markdown file')
    parser.add_argument('-o', '--output', help='Output file (default: input name with format extension)')
    parser.add_argument('--format', choices=['pdf', 'docx', 'xlsx', 'pptx'], default='pdf',
                       help='Output format (default: pdf)')
    parser.add_argument('--append', metavar='FILE', 
                       help='Append to existing file (PDF or PPTX only)')
    parser.add_argument('--replace', metavar='FILE',
                       help='Replace matching content in existing file (PDF or PPTX only)')
    
    args = parser.parse_args()
    
    # Validate input
    if not os.path.exists(args.markdown_file):
        print(f"Error: Markdown file not found: {args.markdown_file}")
        sys.exit(1)
    
    # Check format availability
    if args.format == 'pdf' and not PDF_AVAILABLE:
        print("Error: PDF support not available. Install: pip install xhtml2pdf pypdf pdfplumber")
        sys.exit(1)
    elif args.format == 'docx' and not DOCX_AVAILABLE:
        print("Error: DOCX support not available. Install: pip install python-docx")
        sys.exit(1)
    elif args.format == 'xlsx' and not XLSX_AVAILABLE:
        print("Error: XLSX support not available. Install: pip install openpyxl")
        sys.exit(1)
    elif args.format == 'pptx' and not PPTX_AVAILABLE:
        print("Error: PPTX support not available. Install: pip install python-pptx")
        sys.exit(1)
    
    # Validate mode combinations
    if args.append and args.replace:
        print("Error: Cannot use both --append and --replace modes")
        sys.exit(1)
    
    if (args.append or args.replace) and args.format not in ['pdf', 'pptx']:
        print(f"Error: --append and --replace are only supported for PDF and PPTX formats")
        sys.exit(1)
    
    # Determine output file
    if args.output:
        output_file = args.output
    else:
        base = os.path.splitext(args.markdown_file)[0]
        output_file = f"{base}.{args.format}"
    
    # Execute conversion
    try:
        with open(args.markdown_file, 'r', encoding='utf-8') as f:
            md_content = f.read()
        
        if args.format == 'pdf':
            if args.append:
                if not os.path.exists(args.append):
                    print(f"Error: File not found: {args.append}")
                    sys.exit(1)
                append_to_pdf(args.markdown_file, args.append, output_file)
            elif args.replace:
                if not os.path.exists(args.replace):
                    print(f"Error: File not found: {args.replace}")
                    sys.exit(1)
                replace_pdf_pages(args.markdown_file, args.replace, output_file)
            else:
                create_pdf(args.markdown_file, output_file)
        
        elif args.format == 'docx':
            markdown_to_docx(md_content, output_file)
            print(f"Word document created successfully: {output_file}")
        
        elif args.format == 'xlsx':
            markdown_to_xlsx(md_content, output_file)
            print(f"Excel spreadsheet created successfully: {output_file}")
        
        elif args.format == 'pptx':
            if args.append:
                if not os.path.exists(args.append):
                    print(f"Error: File not found: {args.append}")
                    sys.exit(1)
                append_to_pptx(args.markdown_file, args.append, output_file)
            elif args.replace:
                if not os.path.exists(args.replace):
                    print(f"Error: File not found: {args.replace}")
                    sys.exit(1)
                replace_pptx_slides(args.markdown_file, args.replace, output_file)
            else:
                markdown_to_pptx(md_content, output_file)
                print(f"PowerPoint presentation created successfully: {output_file}")
    
    except Exception as e:
        print(f"Error: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == '__main__':
    main()
